"""Deterministic, local-worker evidence probes for the benchmark report.

The raster attack matrix in :mod:`bench.runners.run` deliberately does not
pretend to be a detector.  This module supplements it with small, local calls
to the worker's public carrier and format-adapter APIs.  Each probe preserves
the worker's raw evidence and deliberately stops short of attribution: a
benchmark has no server-resolved issuance or web-session binding.

This is not a model benchmark.  The current image detector is the documented
metadata fallback and the current screen detector is candidate correlation.
Their version strings and the installed package versions are written into the
report so a result cannot be misrepresented as a different detector run.
"""

from __future__ import annotations

import io
import re
import sys
from importlib import import_module
from pathlib import Path
from typing import Any
import zipfile

import numpy as np
from PIL import Image

from bench.attacks import Attack, apply_attack


_IMAGE_MIME = "image/png"
_DOCX_MIME = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
_PPTX_MIME = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
_PDF_MIME = "application/pdf"
_BENCHMARK_PROFILE_VERSION = "benchmark-profile-2026-08"
_PRIMARY_HANDLE = "0123456789abcdef0123456789abcdef"
_SECONDARY_HANDLE = "fedcba9876543210fedcba9876543210"


def _worker_root() -> Path:
    return Path(__file__).resolve().parents[1] / "services" / "watermark-worker"


def _worker_api() -> dict[str, Any]:
    """Load the checked-out worker package, never a similarly named installed app.

    The worker intentionally has no separate wheel name for its import package.
    Add its source root ahead of site packages before importing its public APIs.
    Failing loudly when another ``app`` package is already loaded prevents a
    benchmark from accidentally measuring unrelated code.
    """

    root = _worker_root()
    existing = sys.modules.get("app")
    if existing is not None:
        module_path = Path(str(getattr(existing, "__file__", ""))).resolve()
        if root not in module_path.parents:
            raise RuntimeError("cannot run worker benchmark after an unrelated 'app' package was imported")
    elif str(root) not in sys.path:
        sys.path.insert(0, str(root))

    return {
        "Artifact": import_module("app.models").Artifact,
        "CarrierProfile": import_module("app.models").CarrierProfile,
        "TraceIdentity": import_module("app.models").TraceIdentity,
        "ImageCodeCarrier": import_module("app.carriers.image_code").ImageCodeCarrier,
        "ScreenTileCarrier": import_module("app.carriers.screen_tile").ScreenTileCarrier,
        "NativeStructureCarrier": import_module("app.carriers.structure").NativeStructureCarrier,
        "AdapterRegistry": import_module("app.formats.registry").AdapterRegistry,
    }


def _evidence(evidence: Any) -> dict[str, Any]:
    return {
        "carrier": evidence.carrier,
        "detectorVersion": evidence.detector_version,
        "score": round(float(evidence.score), 8),
        "raw": evidence.raw,
        "warnings": list(evidence.warnings),
    }


def _result(
    scenario: str,
    *,
    channel: str,
    corpus: str,
    evidence: Any | None,
    status: str,
    reason: str,
    extra_raw: dict[str, Any] | None = None,
) -> dict[str, Any]:
    raw = {"carrierEvidence": _evidence(evidence) if evidence is not None else None}
    if extra_raw:
        raw.update(extra_raw)
    return {
        "scenario": scenario,
        "channel": channel,
        "corpus": corpus,
        "attribution": {"status": status, "candidate": None, "reason": reason},
        "rawDetectorEvidence": raw,
    }


def _docx_source() -> bytes:
    from docx import Document

    document = Document()
    document.add_paragraph("TraceAnytong deterministic benchmark document")
    out = io.BytesIO()
    document.save(out)
    return out.getvalue()


def _pptx_source() -> bytes:
    from pptx import Presentation

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.shapes.add_textbox(0, 0, 1_000_000, 1_000_000).text_frame.text = "TraceAnytong deterministic benchmark slide"
    out = io.BytesIO()
    presentation.save(out)
    return out.getvalue()


def _normalized_openxml(data: bytes) -> bytes:
    """Freeze ZIP metadata for a deterministic structural-evidence fixture.

    Office libraries otherwise write the current DOS timestamp into each ZIP
    entry. That does not change the native carrier evidence, but it would change
    the source fingerprint recorded by the structure detector on each run.
    """

    source = zipfile.ZipFile(io.BytesIO(data))
    try:
        out = io.BytesIO()
        with zipfile.ZipFile(out, "w", compression=zipfile.ZIP_DEFLATED, compresslevel=9) as target:
            for name in sorted(source.namelist()):
                entry = zipfile.ZipInfo(name, date_time=(2020, 1, 1, 0, 0, 0))
                entry.compress_type = zipfile.ZIP_DEFLATED
                entry.external_attr = 0o600 << 16
                target.writestr(entry, source.read(name), compress_type=zipfile.ZIP_DEFLATED, compresslevel=9)
        return out.getvalue()
    finally:
        source.close()


def _minimal_pdf(*, marker: str | None = None) -> bytes:
    """Create a fixed PDF fixture without a time-varying producer timestamp."""

    contents = b"BT /F1 12 Tf 72 720 Td (TraceAnytong deterministic benchmark PDF) Tj ET"
    info = b""
    if marker:
        info = f" /Keywords ({marker})".encode("ascii")
    objects = [
        b"<< /Type /Catalog /Pages 2 0 R >>",
        b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>",
        b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] /Resources << /Font << /F1 4 0 R >> >> /Contents 5 0 R >>",
        b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>",
        b"<< /Length " + str(len(contents)).encode("ascii") + b" >>\nstream\n" + contents + b"\nendstream",
        b"<<" + info + b" >>",
    ]
    out = bytearray(b"%PDF-1.4\n%\xe2\xe3\xcf\xd3\n")
    offsets = [0]
    for index, body in enumerate(objects, start=1):
        offsets.append(len(out))
        out.extend(f"{index} 0 obj\n".encode("ascii"))
        out.extend(body)
        out.extend(b"\nendobj\n")
    xref = len(out)
    out.extend(f"xref\n0 {len(objects) + 1}\n".encode("ascii"))
    out.extend(b"0000000000 65535 f \n")
    out.extend(b"".join(f"{offset:010d} 00000 n \n".encode("ascii") for offset in offsets[1:]))
    out.extend(f"trailer\n<< /Size {len(objects) + 1} /Root 1 0 R /Info 6 0 R >>\nstartxref\n{xref}\n%%EOF\n".encode("ascii"))
    return bytes(out)


def _canonical_pdf_for_structure_probe(artifact: Any, Artifact: Any) -> Any:
    """Remove the random PDF trailer ID before hashing deterministic fixture evidence.

    PyMuPDF assigns a fresh trailer ``/ID`` when the PDF adapter writes a
    personalized artifact.  It is not carrier evidence, but the structure
    detector correctly includes the bytes' SHA-256 in its raw output.  The
    benchmark therefore removes only that documented random field before the
    local structure probe and records the normalization alongside the result.
    This helper is fixture-only; it must not be used for uploaded evidence.
    """

    canonical = re.sub(rb"/ID\s*\[\s*<[^>]+>\s*<[^>]+>\s*\]", b"", artifact.data)
    if canonical == artifact.data:
        raise ValueError("expected personalized PDF fixture to contain a trailer ID")
    return Artifact(canonical, artifact.mime_type, artifact.filename)


def _screen_raster(carrier: Any, identities: tuple[Any, ...], profile: Any) -> Image.Image:
    """Build a repeatable low-energy rendered-screen fixture from worker fields."""

    fields = [carrier.field(identity, profile) for identity in identities]
    base = np.full((256, 320), 185.0, dtype=np.float32)
    for field in fields:
        base += 16 * np.tile(field, (4, 5))[:, :320]
    return Image.fromarray(np.clip(base, 0, 255).astype(np.uint8), "L")


def _screen_candidate_scores(first: Any, second: Any, *, transform: dict[str, Any] | None = None) -> dict[str, Any]:
    """Keep both candidates and transform metadata in every screen probe.

    These labels are fixed benchmark aliases, not trace identities.  The raw
    detector records let a reader inspect candidate separation without the
    benchmark claiming that either candidate was server-resolved.
    """

    raw: dict[str, Any] = {
        "candidateScores": [
            {"candidate": "primary", **_evidence(first)},
            {"candidate": "secondary", **_evidence(second)},
        ],
        "crossCandidateMargin": round(float(first.score - second.score), 8),
    }
    if transform is not None:
        raw["transform"] = transform
    return raw


def _versions(api: dict[str, Any]) -> dict[str, str]:
    import fitz
    import docx
    import pptx
    import PIL

    return {
        "python": ".".join(map(str, sys.version_info[:3])),
        "numpy": np.__version__,
        "pillow": PIL.__version__,
        "pymupdf": str(fitz.VersionBind),
        "python-docx": str(docx.__version__),
        "python-pptx": str(pptx.__version__),
        "imageCarrier": api["ImageCodeCarrier"].carrier_version,
        "imageDetector": api["ImageCodeCarrier"].detector_version,
        "screenDetector": api["ScreenTileCarrier"].detector_version,
        "structureCarrier": api["NativeStructureCarrier"].carrier_version,
        "structureDetector": api["NativeStructureCarrier"].detector_version,
    }


def collect_worker_evidence(fixtures_root: Path) -> dict[str, Any]:
    """Run local worker detectors against deterministic positive, negative and ambiguous probes."""

    api = _worker_api()
    Artifact = api["Artifact"]
    CarrierProfile = api["CarrierProfile"]
    TraceIdentity = api["TraceIdentity"]
    profile = CarrierProfile(
        "benchmark-profile",
        _BENCHMARK_PROFILE_VERSION,
        "benchmark-key-v1",
        b"traceanytong-benchmark-key-material-v1",
        strength=0.08,
        tile_size=64,
    )
    primary = TraceIdentity(_PRIMARY_HANDLE, "web_session", _BENCHMARK_PROFILE_VERSION, 1_725_000_000)
    secondary = TraceIdentity(_SECONDARY_HANDLE, "web_session", _BENCHMARK_PROFILE_VERSION, 1_725_000_001)
    results: list[dict[str, Any]] = []

    image_carrier = api["ImageCodeCarrier"]()
    image_source = (fixtures_root / "image-positive-01.png").read_bytes()
    personalized_image = image_carrier.embed(
        Artifact(image_source, _IMAGE_MIME, "image-positive-01.png"),
        TraceIdentity(_PRIMARY_HANDLE, "issuance", _BENCHMARK_PROFILE_VERSION, 1_725_000_000),
        profile,
        wm_code=0xAABBCCDD,
    )
    recovered = image_carrier.detect(personalized_image.artifact, profile)
    results.append(_result(
        "image-native-metadata-recovery",
        channel="image",
        corpus="acceptance",
        evidence=recovered,
        status="UNMEASURED",
        reason="A benchmark has no server-resolved wmCode-to-issuance binding, so recovered code evidence is not attribution.",
    ))
    decoded = Image.open(io.BytesIO(personalized_image.artifact.data)).convert("RGB")
    transcoded = io.BytesIO()
    decoded.save(transcoded, format="PNG", optimize=False)
    raster_recovered = image_carrier.detect(Artifact(transcoded.getvalue(), _IMAGE_MIME, "metadata-stripped.png"), profile)
    results.append(_result(
        "image-raster-metadata-stripped",
        channel="image",
        corpus="acceptance",
        evidence=raster_recovered,
        status="UNMEASURED",
        reason="The repeated visual code survives metadata stripping, but a benchmark has no server-resolved wmCode-to-issuance binding.",
    ))
    jpeg = io.BytesIO()
    decoded.save(jpeg, format="JPEG", quality=60, optimize=False)
    jpeg_recovered = image_carrier.detect(Artifact(jpeg.getvalue(), "image/jpeg", "metadata-stripped-q60.jpg"), profile)
    results.append(_result(
        "image-raster-jpeg-60",
        channel="image",
        corpus="acceptance",
        evidence=jpeg_recovered,
        status="UNMEASURED",
        reason="Controlled JPEG quality 60 retains measured visual-code evidence, but a benchmark has no server-resolved wmCode-to-issuance binding.",
        extra_raw={"transform": {"name": "jpeg", "value": 60, "metadataStripped": True}},
    ))
    resized = decoded.resize((480, 300), Image.Resampling.LANCZOS)
    resized_stream = io.BytesIO()
    resized.save(resized_stream, format="PNG", optimize=False)
    resized_recovered = image_carrier.detect(Artifact(resized_stream.getvalue(), _IMAGE_MIME, "metadata-stripped-resize-075.png"), profile)
    results.append(_result(
        "image-raster-resize-0.75",
        channel="image",
        corpus="acceptance",
        evidence=resized_recovered,
        status="UNMEASURED",
        reason="Controlled 0.75 resize retains measured visual-code evidence, but a benchmark has no server-resolved wmCode-to-issuance binding.",
        extra_raw={"transform": {"name": "resize", "value": 0.75, "metadataStripped": True}},
    ))
    cropped = decoded.crop((160, 100, 480, 300))
    cropped_stream = io.BytesIO()
    cropped.save(cropped_stream, format="PNG", optimize=False)
    crop_insufficient = image_carrier.detect(Artifact(cropped_stream.getvalue(), _IMAGE_MIME, "metadata-stripped-crop-05.png"), profile)
    results.append(_result(
        "image-raster-crop-0.5-limit",
        channel="image",
        corpus="acceptance",
        evidence=crop_insufficient,
        status="INSUFFICIENT",
        reason="A centered 50% crop of this 640x400 fixture cannot retain all repeated-grid symbols; this fallback reports insufficiency rather than guessing a wmCode.",
        extra_raw={"transform": {"name": "crop", "value": 0.5, "metadataStripped": True}},
    ))
    unmarked = Image.open(fixtures_root / "negative-image-01.png").convert("RGB")
    unmarked_stream = io.BytesIO()
    unmarked.save(unmarked_stream, format="PNG", optimize=False)
    unmarked_evidence = image_carrier.detect(Artifact(unmarked_stream.getvalue(), _IMAGE_MIME, "unwatermarked.png"), profile)
    results.append(_result(
        "image-unwatermarked-negative",
        channel="image",
        corpus="negative",
        evidence=unmarked_evidence,
        status="INSUFFICIENT",
        reason="An unwatermarked deterministic fixture has no CRC-valid, threshold-supported visual code.",
    ))

    screen_carrier = api["ScreenTileCarrier"]()
    matched_screen = _screen_raster(screen_carrier, (primary,), profile)
    matched = screen_carrier.detect_candidate(matched_screen, primary, profile)
    mismatched = screen_carrier.detect_candidate(matched_screen, secondary, profile)
    results.append(_result(
        "screen-candidate-matched",
        channel="screen",
        corpus="acceptance",
        evidence=matched,
        status="UNMEASURED",
        reason="Candidate correlation is raw support only until an authorized server resolves a session binding.",
        extra_raw=_screen_candidate_scores(matched, mismatched),
    ))
    for attack, status, reason in (
        (
            Attack("jpeg", 60),
            "UNMEASURED",
            "JPEG recompression retains a measured candidate separation, but a benchmark has no authorized server-side session binding.",
        ),
        (
            Attack("crop", 0.5),
            "UNMEASURED",
            "The retained crop has a measured candidate separation, but a benchmark has no authorized server-side session binding.",
        ),
        (
            Attack("resize", 0.75),
            "INSUFFICIENT",
            "The current detector has no scale normalization for this resize; the measured scores do not separate the expected candidate and cannot attribute.",
        ),
    ):
        transformed = apply_attack(matched_screen, attack)
        transformed_primary = screen_carrier.detect_candidate(transformed, primary, profile)
        transformed_secondary = screen_carrier.detect_candidate(transformed, secondary, profile)
        results.append(_result(
            f"screen-candidate-{attack.name}-{attack.value:g}",
            channel="screen",
            # This is still an acceptance-corpus input.  Its INSUFFICIENT
            # outcome is a measured detector limitation, not an ambiguous
            # source fixture relabeled to make the matrix look healthier.
            corpus="acceptance",
            evidence=transformed_primary,
            status=status,
            reason=reason,
            extra_raw=_screen_candidate_scores(
                transformed_primary,
                transformed_secondary,
                transform={
                    "name": attack.name,
                    "value": attack.value,
                    "sourceGeometry": {"width": matched_screen.width, "height": matched_screen.height},
                    "artifactGeometry": {"width": transformed.width, "height": transformed.height},
                },
            ),
        ))
    ambiguous_screen = _screen_raster(screen_carrier, (primary, secondary), profile)
    first = screen_carrier.detect_candidate(ambiguous_screen, primary, profile)
    second = screen_carrier.detect_candidate(ambiguous_screen, secondary, profile)
    results.append(_result(
        "screen-two-candidate-ambiguous",
        channel="screen",
        corpus="ambiguous",
        evidence=first if first.score >= second.score else second,
        status="INSUFFICIENT",
        reason="Two candidate fields are deliberately combined; close cross-candidate correlation must never be attributed.",
        extra_raw={
            **_screen_candidate_scores(first, second),
            "crossCandidateMargin": round(float(abs(first.score - second.score)), 8),
        },
    ))
    blank = Image.new("L", (320, 256), 185)
    blank_evidence = screen_carrier.detect_candidate(blank, primary, profile)
    results.append(_result(
        "screen-blank-negative",
        channel="screen",
        corpus="negative",
        evidence=blank_evidence,
        status="INSUFFICIENT",
        reason="No carrier field is present in a flat capture.",
    ))

    registry = api["AdapterRegistry"]()
    structure = api["NativeStructureCarrier"]()
    for label, mime, source in (
        ("docx", _DOCX_MIME, _normalized_openxml(_docx_source())),
        ("pptx", _PPTX_MIME, _normalized_openxml(_pptx_source())),
    ):
        source_artifact = Artifact(source, mime, f"benchmark.{label}")
        personalized = registry.for_mime(mime).personalize(source_artifact, primary, profile)
        # The adapter can restore ZIP timestamps while saving. Freeze them again
        # before fingerprinted extraction so report determinism measures carrier
        # semantics rather than archive clock metadata.
        normalized_personalized = Artifact(_normalized_openxml(personalized.artifact.data), mime, personalized.artifact.filename)
        evidence = structure.detect_candidate(normalized_personalized, primary, profile)
        results.append(_result(
            f"native-{label}-structure-support",
            channel="structure",
            corpus="acceptance",
            evidence=evidence,
            status="INSUFFICIENT",
            reason="Native markers and carrier-shaped placement are provenance support, explicitly not attribution.",
        ))
        unmarked = structure.detect_candidate(source_artifact, primary, profile)
        results.append(_result(
            f"native-{label}-unmarked-negative",
            channel="structure",
            corpus="negative",
            evidence=unmarked,
            status="INSUFFICIENT",
            reason="An unmarked native artifact has no protocol-shaped marker or candidate match.",
        ))

    marker = f"TraceAnytong:{_PRIMARY_HANDLE};profile:{_BENCHMARK_PROFILE_VERSION}"
    personalized_pdf = registry.for_mime(_PDF_MIME).personalize(
        Artifact(_minimal_pdf(), _PDF_MIME, "benchmark.pdf"), primary, profile
    )
    canonical_pdf = _canonical_pdf_for_structure_probe(personalized_pdf.artifact, Artifact)
    personalized_pdf_evidence = structure.detect_candidate(canonical_pdf, primary, profile)
    results.append(_result(
        "native-pdf-structure-support",
        channel="structure",
        corpus="acceptance",
        evidence=personalized_pdf_evidence,
        status="INSUFFICIENT",
        reason="Native PDF markers and measured placements are provenance support, explicitly not attribution.",
        extra_raw={
            "benchmarkNormalization": {
                "removedPdfTrailerId": True,
                "reason": "PyMuPDF creates a random trailer ID; it is removed only to make the generated fixture's raw structure hash reproducible.",
            },
        },
    ))
    pdf_evidence = structure.detect_candidate(Artifact(_minimal_pdf(marker=marker), _PDF_MIME, "marker-only.pdf"), primary, profile)
    results.append(_result(
        "native-pdf-marker-only",
        channel="structure",
        corpus="ambiguous",
        evidence=pdf_evidence,
        status="INSUFFICIENT",
        reason="A PDF metadata marker without a measured carrier-shaped placement is explicitly qualified and cannot attribute.",
    ))

    results.sort(key=lambda item: str(item["scenario"]))
    statuses = {item["attribution"]["status"] for item in results}
    return {
        "schemaVersion": "0.1",
        "deterministic": True,
        "versions": _versions(api),
        "results": results,
        "decisionCounts": {status: sum(item["attribution"]["status"] == status for item in results) for status in sorted(statuses)},
        "confirmedAttributions": 0,
    }
