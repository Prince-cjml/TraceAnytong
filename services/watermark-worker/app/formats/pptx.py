from __future__ import annotations

import io

from pptx import Presentation

from ..carriers.screen_tile import ScreenTileCarrier
from ..errors import InvalidArtifactError
from ..fingerprint.perceptual import PerceptualFingerprinter
from ..models import Artifact, CarrierEvidence, CarrierProfile, PersonalizationResult, TraceIdentity
from .base import BaseAdapter


class PptxFormatAdapter(BaseAdapter):
    mime_types = frozenset({"application/vnd.openxmlformats-officedocument.presentationml.presentation"})

    def __init__(self, carrier: ScreenTileCarrier | None = None) -> None:
        self.carrier = carrier or ScreenTileCarrier()

    def personalize(self, artifact: Artifact, trace_identity: TraceIdentity, profile: CarrierProfile, **_: object) -> PersonalizationResult:
        self.require_support(artifact.mime_type)
        try:
            presentation = Presentation(io.BytesIO(artifact.data))
            tile_file = io.BytesIO()
            self.carrier.tile_rgba(trace_identity, profile).save(tile_file, "PNG")
            for slide in presentation.slides:
                # Raster alpha keeps this low-visible while all original slide objects stay editable.
                slide.shapes.add_picture(io.BytesIO(tile_file.getvalue()), 0, 0, width=presentation.slide_width, height=presentation.slide_height)
            presentation.core_properties.comments = f"TraceAnytong opaque trace {trace_identity.trace_handle}; profile {profile.profile_version}"
            out = io.BytesIO()
            presentation.save(out)
        except Exception as exc:
            raise InvalidArtifactError("PPTX bytes cannot be personalized") from exc
        result = Artifact(out.getvalue(), artifact.mime_type, artifact.filename)
        return PersonalizationResult(result, CarrierEvidence("screen", self.carrier.detector_version, 1.0, {"slides": len(presentation.slides), "placement": "full-slide-alpha-picture", "carrierVersion": profile.carrier_version, "provenance": "core-properties-comments"}), PerceptualFingerprinter().index(result))

    def render_reference(self, artifact: Artifact) -> list[Artifact]:
        self.require_support(artifact.mime_type)
        return []
