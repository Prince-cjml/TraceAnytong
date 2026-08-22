from __future__ import annotations

import io
from dataclasses import replace

import fitz
import pytest
from docx import Document
from pptx import Presentation

from app.carriers.structure import NativeStructureCarrier
from app.errors import InvalidArtifactError, UnsupportedFormatError
from app.formats.registry import AdapterRegistry
from app.models import Artifact


PDF = "application/pdf"
DOCX = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
PPTX = "application/vnd.openxmlformats-officedocument.presentationml.presentation"


def _pdf_artifact() -> Artifact:
    document = fitz.open()
    page = document.new_page(width=612, height=792)
    page.insert_text((72, 72), "A deterministic PDF evidence fixture")
    data = document.tobytes()
    document.close()
    return Artifact(data, PDF, "fixture.pdf")


def _docx_artifact() -> Artifact:
    document = Document()
    document.add_paragraph("A deterministic DOCX evidence fixture")
    stream = io.BytesIO()
    document.save(stream)
    return Artifact(stream.getvalue(), DOCX, "fixture.docx")


def _pptx_artifact() -> Artifact:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.shapes.add_textbox(0, 0, 1000000, 1000000).text_frame.text = "A deterministic PPTX evidence fixture"
    stream = io.BytesIO()
    presentation.save(stream)
    return Artifact(stream.getvalue(), PPTX, "fixture.pptx")


@pytest.mark.parametrize("source_factory", [_pdf_artifact, _docx_artifact, _pptx_artifact])
def test_native_structure_evidence_reports_marker_and_carrier_placement(source_factory, identity, profile):
    source = source_factory()
    personalized = AdapterRegistry().for_mime(source.mime_type).personalize(source, identity, profile)

    evidence = NativeStructureCarrier().detect_candidate(personalized.artifact, identity, profile)

    assert evidence.carrier == "structure"
    assert evidence.detector_version == "native-structure-detector-v1"
    assert evidence.raw["sourceSha256"]
    assert evidence.raw["tools"]["python"]
    assert evidence.raw["provenance"]["markers"] == [{"traceHandle": identity.trace_handle, "profileVersion": profile.profile_version}]
    assert evidence.raw["candidateMatches"] == [{
        "traceHandle": identity.trace_handle,
        "scope": identity.scope,
        "createdAt": identity.created_at,
        "markerProfileVersion": profile.profile_version,
        "profileVersionMatches": True,
    }]
    assert evidence.raw["scoreComponents"]["nativeMarker"] == 0.45
    assert evidence.raw["scoreMeaning"].endswith("not attribution")
    assert evidence.score >= 0.8


def test_unmarked_docx_preserves_structure_but_has_no_candidate_match(identity, profile):
    evidence = NativeStructureCarrier().detect_candidate(_docx_artifact(), identity, profile)

    assert evidence.raw["provenance"]["markers"] == []
    assert evidence.raw["candidateMatches"] == []
    assert evidence.raw["scoreComponents"]["candidateMarkerMatch"] == 0.0
    assert evidence.score == 0.0
    assert any("no protocol-shaped" in warning for warning in evidence.warnings)


def test_pdf_marker_without_carrier_placement_is_explicitly_qualified(identity, profile):
    document = fitz.open()
    document.new_page()
    metadata = document.metadata
    metadata["keywords"] = f"TraceAnytong:{identity.trace_handle};profile:{profile.profile_version}"
    document.set_metadata(metadata)
    artifact = Artifact(document.tobytes(), PDF, "marker-only.pdf")
    document.close()

    evidence = NativeStructureCarrier().detect_candidate(artifact, identity, profile)

    assert evidence.score == pytest.approx(0.65)
    assert evidence.raw["pdf"]["visibleImagePlacementCount"] == 0
    assert any("without a measured carrier-shaped placement" in warning for warning in evidence.warnings)


def test_profile_mismatched_candidate_is_exposed_but_never_counts_as_a_match(identity, profile):
    source = _pptx_artifact()
    personalized = AdapterRegistry().for_mime(source.mime_type).personalize(source, identity, profile)
    mismatched = replace(identity, profile_version="other-profile-v1")

    evidence = NativeStructureCarrier().detect_candidate(personalized.artifact, mismatched, profile)

    assert evidence.raw["candidateMatches"][0]["profileVersionMatches"] is False
    assert evidence.raw["scoreComponents"]["candidateMarkerMatch"] == 0.0
    assert evidence.score == pytest.approx(0.8)


def test_pdf_metadata_is_inspected_for_opaque_markers_but_arbitrary_values_are_not_returned(profile):
    document = fitz.open()
    document.new_page()
    metadata = document.metadata
    metadata["author"] = "private author value"
    document.set_metadata(metadata)
    artifact = Artifact(document.tobytes(), PDF, "author.pdf")
    document.close()

    evidence = NativeStructureCarrier().detect(artifact, profile)

    assert "author" in evidence.raw["pdf"]["metadataFieldsPresent"]
    assert "private author value" not in repr(evidence.raw)


def test_invalid_and_unsupported_artifacts_are_typed(profile):
    detector = NativeStructureCarrier()
    with pytest.raises(UnsupportedFormatError):
        detector.detect(Artifact(b"plain", "text/plain"), profile)
    with pytest.raises(InvalidArtifactError):
        detector.detect(Artifact(b"not-a-pdf", PDF), profile)
