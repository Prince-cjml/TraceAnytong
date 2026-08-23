import base64
import hashlib
import io

import fitz
import pytest
from PIL import Image, PngImagePlugin
from docx import Document
from pptx import Presentation

from app.control_plane import ClaimedJob
from app.errors import InputDownloadError, LeaseLostError
from app.execution import JobRunner, WorkerSettings
from app.carriers.image_code import ImageCodeCarrier
from app.carriers.screen_tile import ScreenTileCarrier
from app.carriers.structure import NativeStructureCarrier
from app.fingerprint.perceptual import PerceptualFingerprinter
from app.formats.registry import AdapterRegistry
from app.formats.office_renderer import OfficeRenderResult, OfficeRenderer
from app.models import Artifact, CarrierEvidence, CarrierProfile, TraceIdentity


def png_bytes() -> bytes:
    image = Image.new("RGB", (160, 100), (120, 150, 180))
    out = io.BytesIO()
    image.save(out, "PNG")
    return out.getvalue()


def watermarked_png(wm_code: int) -> bytes:
    image = Image.new("RGB", (160, 100), (120, 150, 180))
    out = io.BytesIO()
    metadata = PngImagePlugin.PngInfo()
    metadata.add_text("TraceAnytong-wmCode", str(wm_code))
    image.save(out, "PNG", pnginfo=metadata)
    return out.getvalue()


def visual_watermarked_png(wm_code: int) -> bytes:
    """A metadata-bearing image whose repeated visual code survives fixture transforms."""
    image = Image.new("RGB", (512, 512), (120, 150, 180))
    out = io.BytesIO()
    image.save(out, "PNG")
    profile = CarrierProfile("image-v1", "profile-2026-08", "key-1", b"deterministic-worker-key", tile_size=64)
    identity = TraceIdentity("0123456789abcdef0123456789abcdef", "issuance", "profile-2026-08", 1_725_000_000)
    return ImageCodeCarrier().embed(Artifact(out.getvalue(), "image/png", "visual.png"), identity, profile, wm_code=wm_code).artifact.data


def metadata_stripped_jpeg_resize(data: bytes) -> bytes:
    image = Image.open(io.BytesIO(data)).convert("RGB").resize((384, 384), Image.Resampling.LANCZOS)
    out = io.BytesIO()
    image.save(out, "JPEG", quality=60, optimize=False)
    return out.getvalue()


def structured_docx(trace_handle: str, profile_version: str) -> bytes:
    document = Document()
    document.add_paragraph("Trace evidence fixture")
    source = io.BytesIO()
    document.save(source)
    profile = CarrierProfile("structure-v1", profile_version, "key-1", b"deterministic-structure-key", tile_size=64)
    identity = TraceIdentity(trace_handle, "issuance", profile_version, 1_725_000_000)
    personalized = AdapterRegistry().for_mime("application/vnd.openxmlformats-officedocument.wordprocessingml.document").personalize(
        Artifact(source.getvalue(), "application/vnd.openxmlformats-officedocument.wordprocessingml.document", "fixture.docx"),
        identity,
        profile,
    )
    return personalized.artifact.data


def screen_capture(trace_handle: str, profile_version: str, scope: str = "web_session") -> bytes:
    profile = CarrierProfile("screen-v1", profile_version, "key-1", b"deterministic-screen-key", tile_size=64)
    identity = TraceIdentity(trace_handle, scope, profile_version, 1_725_000_000)
    tile = ScreenTileCarrier().tile_rgba(identity, profile)
    capture = Image.new("RGBA", (256, 256), (235, 240, 250, 255))
    for y in range(0, 256, 64):
        for x in range(0, 256, 64):
            capture.alpha_composite(tile, (x, y))
    out = io.BytesIO()
    capture.convert("RGB").save(out, "PNG")
    return out.getvalue()


def screen_pdf(trace_handle: str, profile_version: str) -> bytes:
    """Create a deterministic PDF issued with the normal screen adapter."""
    document = fitz.open()
    page = document.new_page(width=612, height=792)
    page.insert_text((72, 72), "Screen-profile PDF trace fixture")
    source = document.tobytes()
    document.close()
    profile = CarrierProfile("screen-v1", profile_version, "key-1", b"deterministic-screen-key", tile_size=64)
    identity = TraceIdentity(trace_handle, "issuance", profile_version, 1_725_000_000)
    return AdapterRegistry().for_mime("application/pdf").personalize(
        Artifact(source, "application/pdf", "fixture.pdf"), identity, profile
    ).artifact.data


def screen_native_document(mime: str, trace_handle: str, profile_version: str) -> bytes:
    """Create a native Office artifact without claiming an Office render path."""
    source = io.BytesIO()
    if mime == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
        document = Document()
        document.add_paragraph("Screen-profile DOCX trace fixture")
        document.save(source)
        filename = "fixture.docx"
    elif mime == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
        presentation = Presentation()
        presentation.slides.add_slide(presentation.slide_layouts[6])
        presentation.save(source)
        filename = "fixture.pptx"
    else:
        raise AssertionError(f"unexpected native Office MIME: {mime}")
    profile = CarrierProfile("screen-v1", profile_version, "key-1", b"deterministic-screen-key", tile_size=64)
    identity = TraceIdentity(trace_handle, "issuance", profile_version, 1_725_000_000)
    return AdapterRegistry().for_mime(mime).personalize(
        Artifact(source.getvalue(), mime, filename), identity, profile
    ).artifact.data


def environment() -> dict[str, str]:
    return {
        "WORKER_CONVEX_URL": "https://joyous-anaconda-773.convex.cloud",
        "WORKER_CONVEX_TOKEN": "worker-only-token",
        "WORKER_ID": "test-worker",
        "WORKER_PROFILE_IMAGE_V1_SECRET_BASE64": base64.b64encode(b"deterministic-worker-key").decode(),
        "WORKER_PROFILE_IMAGE_V1_VERSION": "profile-2026-08",
        "WORKER_PROFILE_IMAGE_V1_KEY_VERSION": "key-1",
        "WORKER_PROFILE_IMAGE_V1_TILE_SIZE": "64",
    }


class FakeClient:
    def __init__(self, source: bytes, *, download_error: Exception | None = None, heartbeat_error: Exception | None = None) -> None:
        self.source = source
        self.download_error = download_error
        self.heartbeat_error = heartbeat_error
        self.calls: list[str] = []
        self.completed: dict | None = None
        self.failed: tuple[str, bool] | None = None
        self.trace_candidates: list[dict] = []
        self.completed_cases: list[str] = []
        self.completed_indexes: list[dict] = []
        self.recovered = 0
        self.requeued = 0
        self.job = ClaimedJob("jobs:1", "key", "personalize", "storage:1", "image-v1", 9_999_999)

    def claim(self, worker_id: str, capabilities: list[str]):
        self.calls.append("claim")
        return self.job

    def recover_expired_leases(self) -> int:
        self.calls.append("recover-expired")
        return self.recovered

    def requeue_retries(self) -> int:
        self.calls.append("requeue-retries")
        return self.requeued

    def start(self, worker_id: str, job_id: str) -> None:
        self.calls.append("start")

    def heartbeat(self, worker_id: str, job_id: str) -> None:
        self.calls.append("heartbeat")
        if self.heartbeat_error:
            raise self.heartbeat_error

    def input(self, worker_id: str, job_id: str) -> dict:
        self.calls.append("input")
        return {
            "inputUrl": "https://storage.example/input",
            "mime": "image/png",
            "inputSha256": hashlib.sha256(self.source).hexdigest(),
            "traceHandle": "0123456789abcdef0123456789abcdef",
            "profileVersion": "profile-2026-08",
            "createdAt": 1_725_000_000,
            "wmCode": 42,
            "profileId": "image-v1",
            "profileCarrier": "image",
        }

    def download_input(self, input_url: str) -> bytes:
        self.calls.append("download")
        if self.download_error:
            raise self.download_error
        return self.source

    def create_upload_url(self) -> str:
        self.calls.append("upload-url")
        return "https://storage.example/upload"

    def upload_output(self, upload_url: str, data: bytes, mime_type: str) -> str:
        self.calls.append("upload")
        assert mime_type == "image/png"
        assert data != self.source
        return "storage:derived"

    def complete(self, worker_id: str, job_id: str, output_storage_id: str, output_sha256: str, result: dict) -> dict:
        self.calls.append("complete")
        self.completed = {"storageId": output_storage_id, "sha256": output_sha256, "result": result}
        return {"status": "succeeded"}

    def complete_content_index(self, worker_id: str, job_id: str, result: dict) -> dict:
        self.calls.append("complete-content-index")
        self.completed_indexes.append(result)
        return {"status": "succeeded"}

    def fail(self, worker_id: str, job_id: str, error: str, retryable: bool) -> None:
        self.calls.append("fail")
        self.failed = (error, retryable)

    def record_trace_candidate(self, worker_id: str, job_id: str, args: dict) -> dict:
        self.calls.append("record-candidate")
        assert worker_id == "test-worker"
        assert job_id == self.job.job_id
        self.trace_candidates.append(args)
        return {"candidateId": "candidates:1", "decision": "attributed"}

    def complete_trace_case(self, worker_id: str, job_id: str, case_id: str, failed: bool = False) -> None:
        assert worker_id == "test-worker"
        assert job_id == self.job.job_id
        assert not failed
        self.calls.append("complete-case")
        self.completed_cases.append(case_id)


def runner_for(client: FakeClient, env: dict[str, str] | None = None) -> JobRunner:
    values = env or environment()
    return JobRunner(WorkerSettings.from_env(values), client, env=values)


def test_run_once_personalizes_hashes_direct_uploads_and_keeps_secret_out_of_result() -> None:
    env = environment()
    client = FakeClient(png_bytes())
    outcome = runner_for(client, env).run_once()

    assert outcome.status == "succeeded"
    assert client.calls == ["claim", "start", "input", "heartbeat", "download", "heartbeat", "upload-url", "upload", "heartbeat", "complete"]
    assert client.completed is not None
    assert client.completed["storageId"] == "storage:derived"
    assert len(client.completed["sha256"]) == 64
    result = client.completed["result"]
    assert result["carrierEvidence"]["raw"]["wmCode"] == 42
    assert base64.b64decode(env["WORKER_PROFILE_IMAGE_V1_SECRET_BASE64"]) not in str(result).encode()
    assert "secret" not in result


def test_content_index_runs_without_a_profile_secret_and_uploads_manifest_previews_and_features() -> None:
    client = FakeClient(png_bytes())
    client.job = ClaimedJob("jobs:index", "index-key", "content_index", "storage:source", "source-content-index-v1", 9_999_999)
    def index_input(worker_id: str, job_id: str) -> dict:
        client.calls.append("input")
        return {
            "inputUrl": "https://storage.example/input", "mime": "image/png", "inputSha256": hashlib.sha256(client.source).hexdigest(),
            "versionId": "documentVersions:immutable", "indexVersion": "source-content-index-v1", "maxPages": 200,
        }
    client.input = index_input  # type: ignore[method-assign]
    def index_upload(upload_url: str, data: bytes, mime_type: str) -> str:
        client.calls.append(f"upload:{mime_type}")
        return f"storage:{len(client.calls)}"
    client.upload_output = index_upload  # type: ignore[method-assign]

    outcome = runner_for(client, {key: value for key, value in environment().items() if "PROFILE_" not in key}).run_once()

    assert outcome.status == "succeeded"
    assert client.completed_indexes[0]["status"] == "indexed"
    assert client.completed_indexes[0]["pages"][0]["pageIndex"] == 0
    assert "upload:image/png" in client.calls
    assert client.calls.count("upload:application/json") == 2


def test_content_index_respects_the_leased_page_limit_before_uploading_pages() -> None:
    document = fitz.open()
    for _ in range(2):
        document.new_page(width=72, height=72)
    source = document.tobytes()
    document.close()
    client = FakeClient(source)
    client.job = ClaimedJob("jobs:index-limit", "index-limit-key", "content_index", "storage:source", "source-content-index-v1", 9_999_999)

    def index_input(worker_id: str, job_id: str) -> dict:
        client.calls.append("input")
        return {
            "inputUrl": "https://storage.example/input", "mime": "application/pdf", "inputSha256": hashlib.sha256(client.source).hexdigest(),
            "versionId": "documentVersions:immutable", "indexVersion": "source-content-index-v1", "maxPages": 1,
        }

    client.input = index_input  # type: ignore[method-assign]

    def index_upload(upload_url: str, data: bytes, mime_type: str) -> str:
        client.calls.append(f"upload:{mime_type}")
        return f"storage:{len(client.calls)}"

    client.upload_output = index_upload  # type: ignore[method-assign]
    outcome = runner_for(client, {key: value for key, value in environment().items() if "PROFILE_" not in key}).run_once()

    assert outcome.status == "succeeded"
    assert client.completed_indexes[0]["status"] == "unindexed"
    assert client.completed_indexes[0]["pages"] == []
    assert client.calls.count("upload:image/png") == 0
    assert client.calls.count("upload:application/json") == 1


def test_maintenance_uses_server_controlled_lease_recovery_before_retry_requeue() -> None:
    client = FakeClient(png_bytes())
    client.recovered = 2
    client.requeued = 3

    assert runner_for(client).maintain() == {"recovered": 2, "requeued": 3}
    assert client.calls == ["recover-expired", "requeue-retries"]


def test_download_failure_is_recorded_as_retryable_without_uploading() -> None:
    client = FakeClient(png_bytes(), download_error=InputDownloadError("temporary storage outage"))
    outcome = runner_for(client).run_once()

    assert outcome.status == "failed"
    assert outcome.error_code == "INPUT_DOWNLOAD_ERROR"
    assert client.failed == ("INPUT_DOWNLOAD_ERROR", True)
    assert "upload" not in client.calls


def test_lease_loss_never_attempts_to_complete_or_fail() -> None:
    client = FakeClient(png_bytes(), heartbeat_error=LeaseLostError("job lease is no longer active"))
    outcome = runner_for(client).run_once()

    assert outcome.status == "lease_lost"
    assert client.failed is None
    assert "complete" not in client.calls


def test_profile_configuration_rejects_version_drift_and_keeps_key_internal() -> None:
    env = environment()
    env["WORKER_PROFILE_IMAGE_V1_VERSION"] = "other-profile"
    client = FakeClient(png_bytes())
    outcome = runner_for(client, env).run_once()

    assert outcome.status == "failed"
    assert client.failed == ("PROFILE_CONFIGURATION_ERROR", False)
    assert client.completed is None


def test_integral_convex_json_wm_code_is_normalized_before_personalization() -> None:
    """Convex HTTP encodes its numeric values as Python floats."""
    client = FakeClient(png_bytes())
    original_input = client.input

    def float_code(worker_id: str, job_id: str) -> dict:
        payload = original_input(worker_id, job_id)
        payload["wmCode"] = 42.0
        return payload

    client.input = float_code  # type: ignore[method-assign]
    outcome = runner_for(client).run_once()

    assert outcome.status == "succeeded"
    assert client.completed is not None
    assert client.completed["result"]["carrierEvidence"]["raw"]["wmCode"] == 42


@pytest.mark.parametrize(("mime", "expected_carrier"), [
    ("image/jpeg", "image"),
    ("image/png", "image"),
    ("image/webp", "image"),
    ("application/pdf", "screen"),
    ("application/vnd.openxmlformats-officedocument.wordprocessingml.document", "screen"),
    ("application/vnd.openxmlformats-officedocument.presentationml.presentation", "screen"),
])
def test_personalization_carrier_mapping_matches_implemented_adapters(mime: str, expected_carrier: str) -> None:
    assert JobRunner._personalization_carrier_for_mime(mime) == expected_carrier


def test_personalization_rejects_profile_carrier_incompatible_with_source_mime() -> None:
    client = FakeClient(png_bytes())
    original_input = client.input

    def mismatched_carrier_input(worker_id: str, job_id: str) -> dict:
        payload = original_input(worker_id, job_id)
        payload["profileCarrier"] = "screen"
        return payload

    client.input = mismatched_carrier_input  # type: ignore[method-assign]
    outcome = runner_for(client).run_once()

    assert outcome.status == "failed"
    assert outcome.error_code == "INPUT_INTEGRITY_ERROR"
    assert client.failed == ("INPUT_INTEGRITY_ERROR", False)
    assert "download" not in client.calls


def test_web_tile_job_uploads_a_png_without_downloading_document_bytes() -> None:
    env = environment()
    env.update({
        "WORKER_PROFILE_SCREEN_V1_SECRET_BASE64": base64.b64encode(b"deterministic-screen-key").decode(),
        "WORKER_PROFILE_SCREEN_V1_VERSION": "profile-2026-08",
        "WORKER_PROFILE_SCREEN_V1_TILE_SIZE": "64",
        "WORKER_OFFICE_RENDERER_PATH": "traceanytong-test-no-office-renderer",
    })
    client = FakeClient(png_bytes())
    client.job = ClaimedJob("jobs:tile", "tile-key", "web_tile", None, "screen-v1", 9_999_999)
    original_input = client.input

    def tile_input(worker_id: str, job_id: str) -> dict:
        payload = original_input(worker_id, job_id)
        payload.update({"scope": "web_session", "profileId": "screen-v1", "profileCarrier": "screen", "profileVersion": "profile-2026-08", "wmCode": None})
        return payload

    client.input = tile_input  # type: ignore[method-assign]
    outcome = runner_for(client, env).run_once()

    assert outcome.status == "succeeded"
    assert "download" not in client.calls
    assert client.completed is not None
    assert client.completed["result"]["outputMime"] == "image/png"


def test_trace_job_records_only_a_unique_server_resolved_code_match() -> None:
    env = environment()
    source = watermarked_png(42)
    client = FakeClient(source)
    client.job = ClaimedJob("jobs:trace", "trace-key", "trace", "storage:evidence", "image-v1", 9_999_999, case_id="cases:1")
    original_input = client.input

    def trace_input(worker_id: str, job_id: str) -> dict:
        payload = original_input(worker_id, job_id)
        payload.update({
            "caseId": "cases:1",
            "candidates": [{"issuanceId": "issuances:1", "traceHandle": "0123456789abcdef0123456789abcdef", "wmCode": 42, "outputSha256": hashlib.sha256(source).hexdigest()}],
        })
        return payload

    client.input = trace_input  # type: ignore[method-assign]
    outcome = runner_for(client, env).run_once()

    assert outcome.status == "succeeded"
    assert client.trace_candidates[0]["issuanceId"] == "issuances:1"
    assert client.trace_candidates[0]["fingerprintScore"] == 1.0
    assert client.completed_cases == ["cases:1"]
    assert client.completed is not None
    assert client.completed["storageId"] is None


def test_trace_job_never_requests_attribution_from_metadata_only_recovery() -> None:
    env = environment()
    source = watermarked_png(42)
    client = FakeClient(source)
    client.job = ClaimedJob("jobs:trace", "trace-key", "trace", "storage:evidence", "image-v1", 9_999_999, case_id="cases:1")
    original_input = client.input

    def transformed_trace_input(worker_id: str, job_id: str) -> dict:
        payload = original_input(worker_id, job_id)
        payload.update({
            "caseId": "cases:1",
            "candidates": [{"issuanceId": "issuances:1", "traceHandle": "0123456789abcdef0123456789abcdef", "wmCode": 42, "outputSha256": "0" * 64}],
        })
        return payload

    client.input = transformed_trace_input  # type: ignore[method-assign]
    outcome = runner_for(client, env).run_once()

    assert outcome.status == "succeeded"
    assert client.trace_candidates[0]["fingerprintScore"] == 0.0
    assert client.trace_candidates[0]["requestedDecision"] == "insufficient"
    assert client.trace_candidates[0]["finalConfidence"] == 0.0


def test_trace_job_fuses_visual_code_with_frozen_perceptual_fingerprint_after_transform() -> None:
    source = visual_watermarked_png(42)
    transformed = metadata_stripped_jpeg_resize(source)
    frozen_fingerprint = PerceptualFingerprinter().index(Artifact(source, "image/png", "visual.png"))
    client = FakeClient(transformed)
    client.job = ClaimedJob("jobs:trace-visual-fusion", "trace-visual-fusion-key", "trace", "storage:evidence", "image-v1", 9_999_999, case_id="cases:visual-fusion")
    original_input = client.input

    def visual_fusion_input(worker_id: str, job_id: str) -> dict:
        payload = original_input(worker_id, job_id)
        payload.update({
            "caseId": "cases:visual-fusion", "mime": "image/jpeg",
            "candidates": [{
                "issuanceId": "issuances:visual-fusion", "traceHandle": "0123456789abcdef0123456789abcdef",
                "scope": "issuance", "createdAt": 1_725_000_000, "wmCode": 42,
                "outputSha256": hashlib.sha256(source).hexdigest(), "outputFingerprint": frozen_fingerprint,
            }],
        })
        return payload

    client.input = visual_fusion_input  # type: ignore[method-assign]
    outcome = runner_for(client).run_once()

    assert outcome.status == "succeeded"
    candidate = client.trace_candidates[0]
    assert candidate["requestedDecision"] == "attributed"
    assert candidate["fingerprintScore"] >= JobRunner._MIN_IMAGE_PERCEPTUAL_SCORE
    assert candidate["rawEvidence"]["imageCarrier"]["raw"]["recovery"] == "visual-raster"
    assert candidate["rawEvidence"]["fingerprint"]["fingerprint_version"] == PerceptualFingerprinter.version
    assert candidate["rawEvidence"]["attributionGate"] == {
        "exactOutputSha256": False,
        "visualRasterRecovery": True,
        "minimumPerceptualScore": JobRunner._MIN_IMAGE_PERCEPTUAL_SCORE,
        "perceptualFingerprintSupported": True,
    }


def test_trace_job_refuses_visual_code_when_frozen_perceptual_fingerprint_mismatches() -> None:
    source = visual_watermarked_png(42)
    transformed = metadata_stripped_jpeg_resize(source)
    frozen_fingerprint = PerceptualFingerprinter().index(Artifact(source, "image/png", "visual.png"))
    frozen_fingerprint["dHash"] = f"{int(frozen_fingerprint['dHash'], 16) ^ ((1 << 64) - 1):016x}"
    client = FakeClient(transformed)
    client.job = ClaimedJob("jobs:trace-visual-mismatch", "trace-visual-mismatch-key", "trace", "storage:evidence", "image-v1", 9_999_999, case_id="cases:visual-mismatch")
    original_input = client.input

    def mismatch_input(worker_id: str, job_id: str) -> dict:
        payload = original_input(worker_id, job_id)
        payload.update({
            "caseId": "cases:visual-mismatch", "mime": "image/jpeg",
            "candidates": [{
                "issuanceId": "issuances:visual-mismatch", "traceHandle": "0123456789abcdef0123456789abcdef",
                "scope": "issuance", "createdAt": 1_725_000_000, "wmCode": 42,
                "outputSha256": hashlib.sha256(source).hexdigest(), "outputFingerprint": frozen_fingerprint,
            }],
        })
        return payload

    client.input = mismatch_input  # type: ignore[method-assign]
    outcome = runner_for(client).run_once()

    assert outcome.status == "succeeded"
    candidate = client.trace_candidates[0]
    assert candidate["requestedDecision"] == "insufficient"
    assert candidate["rawEvidence"]["imageCarrier"]["raw"]["recovery"] == "visual-raster"
    assert candidate["rawEvidence"]["attributionGate"]["perceptualFingerprintSupported"] is False


def test_trace_job_records_native_document_structure_as_insufficient_only() -> None:
    trace_handle = "0123456789abcdef0123456789abcdef"
    source = structured_docx(trace_handle, "profile-2026-08")
    env = environment()
    env.update({
        "WORKER_PROFILE_STRUCTURE_V1_SECRET_BASE64": base64.b64encode(b"deterministic-structure-key").decode(),
        "WORKER_PROFILE_STRUCTURE_V1_VERSION": "profile-2026-08",
        "WORKER_PROFILE_STRUCTURE_V1_TILE_SIZE": "64",
    })
    client = FakeClient(source)
    client.job = ClaimedJob("jobs:trace-docx", "trace-docx-key", "trace", "storage:evidence", "structure-v1", 9_999_999, case_id="cases:docx")
    original_input = client.input

    def document_trace_input(worker_id: str, job_id: str) -> dict:
        payload = original_input(worker_id, job_id)
        payload.update({
            "caseId": "cases:docx",
            "mime": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            "profileId": "structure-v1",
            "candidates": [{"issuanceId": "issuances:structure", "traceHandle": trace_handle, "scope": "issuance", "createdAt": 1_725_000_000, "wmCode": None, "outputSha256": hashlib.sha256(source).hexdigest()}],
        })
        return payload

    client.input = document_trace_input  # type: ignore[method-assign]
    outcome = runner_for(client, env).run_once()

    assert outcome.status == "succeeded"
    assert client.trace_candidates[0]["structureScore"] >= 0.8
    assert client.trace_candidates[0]["requestedDecision"] == "insufficient"
    assert client.trace_candidates[0]["rawEvidence"]["nativeStructure"]["raw"]["scoreMeaning"].endswith("not attribution")


def test_trace_job_ranks_candidate_matched_screen_capture() -> None:
    trace_handle = "0123456789abcdef0123456789abcdef"
    source = screen_capture(trace_handle, "profile-2026-08")
    env = environment()
    env.update({
        "WORKER_PROFILE_SCREEN_V1_SECRET_BASE64": base64.b64encode(b"deterministic-screen-key").decode(),
        "WORKER_PROFILE_SCREEN_V1_VERSION": "profile-2026-08",
        "WORKER_PROFILE_SCREEN_V1_TILE_SIZE": "64",
    })
    client = FakeClient(source)
    client.job = ClaimedJob("jobs:trace-screen", "trace-screen-key", "trace", "storage:evidence", "screen-v1", 9_999_999, case_id="cases:screen")
    original_input = client.input

    def screen_trace_input(worker_id: str, job_id: str) -> dict:
        payload = original_input(worker_id, job_id)
        payload.update({
            "caseId": "cases:screen", "profileId": "screen-v1", "profileCarrier": "screen",
            "candidates": [{"webSessionId": "sessions:1", "traceHandle": trace_handle, "scope": "web_session", "createdAt": 1_725_000_000, "wmCode": None, "outputSha256": None}],
        })
        return payload

    client.input = screen_trace_input  # type: ignore[method-assign]
    outcome = runner_for(client, env).run_once()

    assert outcome.status == "succeeded"
    assert client.trace_candidates[0]["webSessionId"] == "sessions:1"
    assert client.trace_candidates[0]["watermarkScore"] >= 0.2
    assert client.trace_candidates[0]["requestedDecision"] == "insufficient"
    assert "candidateScores" in client.trace_candidates[0]["rawEvidence"]


def test_trace_job_records_issuance_scoped_screen_candidate() -> None:
    trace_handle = "0123456789abcdef0123456789abcdef"
    source = screen_capture(trace_handle, "profile-2026-08", "issuance")
    env = environment()
    env.update({
        "WORKER_PROFILE_SCREEN_V1_SECRET_BASE64": base64.b64encode(b"deterministic-screen-key").decode(),
        "WORKER_PROFILE_SCREEN_V1_VERSION": "profile-2026-08",
        "WORKER_PROFILE_SCREEN_V1_TILE_SIZE": "64",
    })
    client = FakeClient(source)
    client.job = ClaimedJob("jobs:trace-screen-issuance", "trace-screen-issuance-key", "trace", "storage:evidence", "screen-v1", 9_999_999, case_id="cases:screen-issuance")
    original_input = client.input

    def screen_trace_input(worker_id: str, job_id: str) -> dict:
        payload = original_input(worker_id, job_id)
        payload.update({
            "caseId": "cases:screen-issuance", "profileId": "screen-v1", "profileCarrier": "screen",
            "candidates": [{"issuanceId": "issuances:screen", "traceHandle": trace_handle, "scope": "issuance", "createdAt": 1_725_000_000, "wmCode": None, "outputSha256": hashlib.sha256(source).hexdigest()}],
        })
        return payload

    client.input = screen_trace_input  # type: ignore[method-assign]
    outcome = runner_for(client, env).run_once()

    assert outcome.status == "succeeded"
    assert client.trace_candidates[0]["issuanceId"] == "issuances:screen"
    assert "webSessionId" not in client.trace_candidates[0]
    assert client.trace_candidates[0]["rawEvidence"]["candidateScores"][0]["scope"] == "issuance"
    assert client.trace_candidates[0]["requestedDecision"] == "insufficient"


def test_trace_job_renders_pdf_screen_evidence_with_page_correlations() -> None:
    trace_handle = "0123456789abcdef0123456789abcdef"
    source = screen_pdf(trace_handle, "profile-2026-08")
    env = environment()
    env.update({
        "WORKER_PROFILE_SCREEN_V1_SECRET_BASE64": base64.b64encode(b"deterministic-screen-key").decode(),
        "WORKER_PROFILE_SCREEN_V1_VERSION": "profile-2026-08",
        "WORKER_PROFILE_SCREEN_V1_TILE_SIZE": "64",
    })
    client = FakeClient(source)
    client.job = ClaimedJob("jobs:trace-pdf-screen", "trace-pdf-screen-key", "trace", "storage:evidence", "screen-v1", 9_999_999, case_id="cases:pdf-screen")
    original_input = client.input

    def pdf_trace_input(worker_id: str, job_id: str) -> dict:
        payload = original_input(worker_id, job_id)
        payload.update({
            "caseId": "cases:pdf-screen", "mime": "application/pdf", "profileId": "screen-v1", "profileCarrier": "screen",
            "candidates": [{"issuanceId": "issuances:pdf", "traceHandle": trace_handle, "scope": "issuance", "createdAt": 1_725_000_000, "wmCode": None, "outputSha256": hashlib.sha256(source).hexdigest()}],
        })
        return payload

    client.input = pdf_trace_input  # type: ignore[method-assign]
    outcome = runner_for(client, env).run_once()

    assert outcome.status == "succeeded"
    assert client.trace_candidates[0]["issuanceId"] == "issuances:pdf"
    assert client.trace_candidates[0]["requestedDecision"] == "insufficient"
    assert client.trace_candidates[0]["rawEvidence"]["pageCorrelations"]
    page = client.trace_candidates[0]["rawEvidence"]["pageCorrelations"][0]
    assert page["page"] == 1
    assert page["screenCorrelation"]["detector_version"] == ScreenTileCarrier.detector_version
    assert client.trace_candidates[0]["rawEvidence"]["screenCorrelation"]["score"] == client.trace_candidates[0]["watermarkScore"]
    assert client.trace_candidates[0]["finalConfidence"] == 0.0
    assert client.trace_candidates[0]["rawEvidence"]["attributionGate"] == {
        "contentMatchAvailable": False,
        "contentMatchStatus": "unavailable",
        "correlationClear": False,
        "warning": JobRunner._SCREEN_CONTENT_MATCH_WARNING,
    }


def test_trace_job_with_pdf_and_empty_snapshot_completes_without_a_match() -> None:
    source = screen_pdf("0123456789abcdef0123456789abcdef", "profile-2026-08")
    env = environment()
    env.update({
        "WORKER_PROFILE_SCREEN_V1_SECRET_BASE64": base64.b64encode(b"deterministic-screen-key").decode(),
        "WORKER_PROFILE_SCREEN_V1_VERSION": "profile-2026-08",
        "WORKER_PROFILE_SCREEN_V1_TILE_SIZE": "64",
    })
    client = FakeClient(source)
    client.job = ClaimedJob("jobs:trace-pdf-empty", "trace-pdf-empty-key", "trace", "storage:evidence", "screen-v1", 9_999_999, case_id="cases:pdf-empty")
    original_input = client.input

    def pdf_trace_input(worker_id: str, job_id: str) -> dict:
        payload = original_input(worker_id, job_id)
        payload.update({
            "caseId": "cases:pdf-empty", "mime": "application/pdf", "profileId": "screen-v1", "profileCarrier": "screen", "candidates": [],
        })
        return payload

    client.input = pdf_trace_input  # type: ignore[method-assign]
    outcome = runner_for(client, env).run_once()

    assert outcome.status == "succeeded"
    assert client.trace_candidates == []
    assert client.completed_cases == ["cases:pdf-empty"]
    assert client.completed is not None
    assert client.completed["result"]["candidateCount"] == 0


@pytest.mark.parametrize("mime", [
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
])
def test_trace_job_uses_structure_only_for_native_office_screen_evidence(mime: str) -> None:
    trace_handle = "0123456789abcdef0123456789abcdef"
    source = screen_native_document(mime, trace_handle, "profile-2026-08")
    env = environment()
    env.update({
        "WORKER_PROFILE_SCREEN_V1_SECRET_BASE64": base64.b64encode(b"deterministic-screen-key").decode(),
        "WORKER_PROFILE_SCREEN_V1_VERSION": "profile-2026-08",
        "WORKER_PROFILE_SCREEN_V1_TILE_SIZE": "64",
        "WORKER_OFFICE_RENDERER_PATH": "traceanytong-test-no-office-renderer",
    })
    client = FakeClient(source)
    client.job = ClaimedJob("jobs:trace-native-screen", "trace-native-screen-key", "trace", "storage:evidence", "screen-v1", 9_999_999, case_id="cases:native-screen")
    original_input = client.input

    def native_trace_input(worker_id: str, job_id: str) -> dict:
        payload = original_input(worker_id, job_id)
        payload.update({
            "caseId": "cases:native-screen", "mime": mime, "profileId": "screen-v1", "profileCarrier": "screen",
            "candidates": [{"issuanceId": "issuances:native", "traceHandle": trace_handle, "scope": "issuance", "createdAt": 1_725_000_000, "wmCode": None, "outputSha256": hashlib.sha256(source).hexdigest()}],
        })
        return payload

    client.input = native_trace_input  # type: ignore[method-assign]
    outcome = runner_for(client, env).run_once()

    assert outcome.status == "succeeded"
    assert client.trace_candidates[0]["requestedDecision"] == "insufficient"
    assert client.trace_candidates[0]["rank"] == 1
    assert client.trace_candidates[0]["watermarkScore"] == 0.0
    assert client.trace_candidates[0]["rawEvidence"]["nativeStructure"]["raw"]["candidateMatches"][0]["traceHandle"] == trace_handle
    rendering = client.trace_candidates[0]["rawEvidence"]["screenVisualCorrelation"]
    assert rendering["attempted"] is False
    assert rendering["available"] is False
    assert rendering["status"] == "unavailable"
    assert rendering["selection"] == "configured"
    assert rendering["externalVersion"] is None
    assert "screenCorrelation" not in client.trace_candidates[0]["rawEvidence"]


@pytest.mark.parametrize("mime", [
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
])
def test_trace_job_uses_optional_office_rendered_pages_for_screen_correlation(monkeypatch, mime: str) -> None:
    trace_handle = "0123456789abcdef0123456789abcdef"
    source = screen_native_document(mime, trace_handle, "profile-2026-08")
    rendered_pdf = screen_pdf(trace_handle, "profile-2026-08")
    env = environment()
    env.update({
        "WORKER_PROFILE_SCREEN_V1_SECRET_BASE64": base64.b64encode(b"deterministic-screen-key").decode(),
        "WORKER_PROFILE_SCREEN_V1_VERSION": "profile-2026-08",
        "WORKER_PROFILE_SCREEN_V1_TILE_SIZE": "64",
    })
    client = FakeClient(source)
    client.job = ClaimedJob("jobs:trace-office-screen", "trace-office-screen-key", "trace", "storage:evidence", "screen-v1", 9_999_999, case_id="cases:office-screen")
    original_input = client.input

    def office_trace_input(worker_id: str, job_id: str) -> dict:
        payload = original_input(worker_id, job_id)
        payload.update({
            "caseId": "cases:office-screen", "mime": mime, "profileId": "screen-v1", "profileCarrier": "screen",
            "candidates": [{"issuanceId": "issuances:office", "traceHandle": trace_handle, "scope": "issuance", "createdAt": 1_725_000_000, "wmCode": None, "outputSha256": hashlib.sha256(source).hexdigest()}],
        })
        return payload

    def fake_render(self, artifact: Artifact) -> OfficeRenderResult:
        assert artifact.data == source
        return OfficeRenderResult(
            Artifact(rendered_pdf, "application/pdf", "fake-rendered.pdf"),
            {"rendererVersion": "libreoffice-pdf-bridge-v1", "available": True, "attempted": True, "status": "rendered", "selection": "configured", "externalVersion": "LibreOffice 24.2.5"},
        )

    monkeypatch.setattr(OfficeRenderer, "render", fake_render)
    client.input = office_trace_input  # type: ignore[method-assign]
    outcome = runner_for(client, env).run_once()

    assert outcome.status == "succeeded"
    assert client.trace_candidates[0]["rawEvidence"]["screenCorrelation"]["detector_version"] == ScreenTileCarrier.detector_version
    assert client.trace_candidates[0]["rawEvidence"]["nativeStructure"]["raw"]["candidateMatches"][0]["traceHandle"] == trace_handle
    renderer = client.trace_candidates[0]["rawEvidence"]["officeRenderer"]
    assert renderer["status"] == "rendered"
    assert renderer["externalVersion"] == "LibreOffice 24.2.5"
    assert renderer["outputPages"] >= 1
    assert client.trace_candidates[0]["requestedDecision"] == "insufficient"
    assert client.trace_candidates[0]["finalConfidence"] == 0.0
    assert client.trace_candidates[0]["rawEvidence"]["attributionGate"]["contentMatchStatus"] == "unavailable"


@pytest.mark.parametrize("mime", [
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
])
def test_trace_job_with_unmatched_native_office_evidence_completes_without_candidates(mime: str) -> None:
    source = screen_native_document(mime, "0123456789abcdef0123456789abcdef", "profile-2026-08")
    env = environment()
    env.update({
        "WORKER_PROFILE_SCREEN_V1_SECRET_BASE64": base64.b64encode(b"deterministic-screen-key").decode(),
        "WORKER_PROFILE_SCREEN_V1_VERSION": "profile-2026-08",
        "WORKER_PROFILE_SCREEN_V1_TILE_SIZE": "64",
        "WORKER_OFFICE_RENDERER_PATH": "traceanytong-test-no-office-renderer",
    })
    client = FakeClient(source)
    client.job = ClaimedJob("jobs:trace-native-empty", "trace-native-empty-key", "trace", "storage:evidence", "screen-v1", 9_999_999, case_id="cases:native-empty")
    original_input = client.input

    def native_trace_input(worker_id: str, job_id: str) -> dict:
        payload = original_input(worker_id, job_id)
        payload.update({
            "caseId": "cases:native-empty", "mime": mime, "profileId": "screen-v1", "profileCarrier": "screen",
            "candidates": [{"issuanceId": "issuances:other", "traceHandle": "fedcba9876543210fedcba9876543210", "scope": "issuance", "createdAt": 1_725_000_001, "wmCode": None, "outputSha256": None}],
        })
        return payload

    client.input = native_trace_input  # type: ignore[method-assign]
    outcome = runner_for(client, env).run_once()

    assert outcome.status == "succeeded"
    assert client.trace_candidates == []
    assert client.completed_cases == ["cases:native-empty"]
    assert client.completed is not None
    assert client.completed["result"]["candidateCount"] == 0
    assert client.completed["result"]["detectorVersion"] == NativeStructureCarrier.detector_version


def test_trace_job_with_empty_screen_snapshot_completes_without_candidates() -> None:
    """An empty immutable snapshot is a no-match result, never a retry loop."""
    env = environment()
    env.update({
        "WORKER_PROFILE_SCREEN_V1_SECRET_BASE64": base64.b64encode(b"deterministic-screen-key").decode(),
        "WORKER_PROFILE_SCREEN_V1_VERSION": "profile-2026-08",
        "WORKER_PROFILE_SCREEN_V1_TILE_SIZE": "64",
    })
    client = FakeClient(png_bytes())
    client.job = ClaimedJob("jobs:trace-screen-empty", "trace-screen-empty-key", "trace", "storage:evidence", "screen-v1", 9_999_999, case_id="cases:screen-empty")
    original_input = client.input

    def screen_trace_input(worker_id: str, job_id: str) -> dict:
        payload = original_input(worker_id, job_id)
        payload.update({
            "caseId": "cases:screen-empty", "profileId": "screen-v1", "profileCarrier": "screen",
            "candidates": [],
        })
        return payload

    client.input = screen_trace_input  # type: ignore[method-assign]
    outcome = runner_for(client, env).run_once()

    assert outcome.status == "succeeded"
    assert client.trace_candidates == []
    assert client.completed_cases == ["cases:screen-empty"]
    assert client.completed is not None
    assert client.completed["result"]["candidateCount"] == 0
    assert client.completed["result"]["detectorVersion"] == ScreenTileCarrier.detector_version


def test_trace_job_retains_the_top_two_ranked_screen_candidates_with_raw_evidence() -> None:
    trace_handle = "0123456789abcdef0123456789abcdef"
    runner_up_handle = "fedcba9876543210fedcba9876543210"
    source = screen_capture(trace_handle, "profile-2026-08")
    env = environment()
    env.update({
        "WORKER_PROFILE_SCREEN_V1_SECRET_BASE64": base64.b64encode(b"deterministic-screen-key").decode(),
        "WORKER_PROFILE_SCREEN_V1_VERSION": "profile-2026-08",
        "WORKER_PROFILE_SCREEN_V1_TILE_SIZE": "64",
    })
    client = FakeClient(source)
    client.job = ClaimedJob("jobs:trace-screen-two", "trace-screen-two-key", "trace", "storage:evidence", "screen-v1", 9_999_999, case_id="cases:screen-two")
    original_input = client.input

    def screen_trace_input(worker_id: str, job_id: str) -> dict:
        payload = original_input(worker_id, job_id)
        payload.update({
            "caseId": "cases:screen-two", "profileId": "screen-v1", "profileCarrier": "screen",
            "candidates": [
                {"webSessionId": "sessions:runner-up", "traceHandle": runner_up_handle, "scope": "web_session", "createdAt": 1_725_000_001, "wmCode": None, "outputSha256": None},
                {"webSessionId": "sessions:top", "traceHandle": trace_handle, "scope": "web_session", "createdAt": 1_725_000_000, "wmCode": None, "outputSha256": None},
            ],
        })
        return payload

    client.input = screen_trace_input  # type: ignore[method-assign]
    outcome = runner_for(client, env).run_once()

    assert outcome.status == "succeeded"
    assert [candidate["rank"] for candidate in client.trace_candidates] == [1, 2]
    assert [candidate["traceHandle"] for candidate in client.trace_candidates] == [trace_handle, runner_up_handle]
    assert client.trace_candidates[0]["watermarkScore"] > client.trace_candidates[1]["watermarkScore"]
    assert client.trace_candidates[0]["watermarkMargin"] == (
        client.trace_candidates[0]["watermarkScore"] - client.trace_candidates[1]["watermarkScore"]
    )
    assert client.trace_candidates[1]["watermarkMargin"] == 0.0
    for candidate in client.trace_candidates:
        raw_evidence = candidate["rawEvidence"]
        assert raw_evidence["candidateRank"] == candidate["rank"]
        assert raw_evidence["screenCorrelation"]["score"] == candidate["watermarkScore"]
        assert [entry["rank"] for entry in raw_evidence["candidateScores"]] == [1, 2]
        assert raw_evidence["candidateScores"][candidate["rank"] - 1]["traceHandle"] == candidate["traceHandle"]
    assert client.trace_candidates[0]["requestedDecision"] == "insufficient"
    assert client.trace_candidates[1]["requestedDecision"] == "insufficient"


def test_trace_job_keeps_a_clear_top_screen_candidate_insufficient_without_content_matching(monkeypatch) -> None:
    top_handle = "0123456789abcdef0123456789abcdef"
    runner_up_handle = "fedcba9876543210fedcba9876543210"
    source = png_bytes()
    env = environment()
    env.update({
        "WORKER_PROFILE_SCREEN_V1_SECRET_BASE64": base64.b64encode(b"deterministic-screen-key").decode(),
        "WORKER_PROFILE_SCREEN_V1_VERSION": "profile-2026-08",
        "WORKER_PROFILE_SCREEN_V1_TILE_SIZE": "64",
    })
    client = FakeClient(source)
    client.job = ClaimedJob("jobs:trace-screen-clear", "trace-screen-clear-key", "trace", "storage:evidence", "screen-v1", 9_999_999, case_id="cases:screen-clear")
    original_input = client.input

    def screen_trace_input(worker_id: str, job_id: str) -> dict:
        payload = original_input(worker_id, job_id)
        payload.update({
            "caseId": "cases:screen-clear", "profileId": "screen-v1", "profileCarrier": "screen",
            "candidates": [
                {"webSessionId": "sessions:runner-up", "traceHandle": runner_up_handle, "scope": "web_session", "createdAt": 1_725_000_001, "wmCode": None, "outputSha256": None},
                {"webSessionId": "sessions:top", "traceHandle": top_handle, "scope": "web_session", "createdAt": 1_725_000_000, "wmCode": None, "outputSha256": None},
            ],
        })
        return payload

    def fixed_candidate_evidence(self, screenshot, identity, profile):
        score = 0.96 if identity.trace_handle == top_handle else 0.30
        return CarrierEvidence(
            carrier="screen", detector_version="screen-correlation-v1", score=score,
            raw={"margin": 0.02, "phase": {"x": 1, "y": 2}, "peak": score, "secondPeak": score - 0.02},
            warnings=(),
        )

    monkeypatch.setattr(ScreenTileCarrier, "detect_candidate", fixed_candidate_evidence)
    client.input = screen_trace_input  # type: ignore[method-assign]
    outcome = runner_for(client, env).run_once()

    assert outcome.status == "succeeded"
    assert [candidate["rank"] for candidate in client.trace_candidates] == [1, 2]
    assert [candidate["requestedDecision"] for candidate in client.trace_candidates] == ["insufficient", "insufficient"]
    assert client.trace_candidates[0]["finalConfidence"] == 0.0
    assert client.trace_candidates[1]["finalConfidence"] == 0.0
    assert client.trace_candidates[0]["watermarkMargin"] == pytest.approx(0.66)
    assert client.trace_candidates[0]["rawEvidence"]["attributionGate"] == {
        "contentMatchAvailable": False,
        "contentMatchStatus": "unavailable",
        "correlationClear": True,
        "warning": JobRunner._SCREEN_CONTENT_MATCH_WARNING,
    }
    assert "immutable content matching" in client.trace_candidates[0]["explanation"]
    assert client.trace_candidates[1]["explanation"].startswith("Runner-up screen candidate")
