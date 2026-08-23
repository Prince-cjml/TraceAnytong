import io

import fitz
import numpy as np
import pytest
from PIL import Image
from docx import Document
from pptx import Presentation

from app.errors import UnsupportedFormatError
from app.formats.registry import AdapterRegistry
from app.models import Artifact


def image_artifact(mime="image/png"):
    image = Image.new("RGB", (160, 100), (120, 150, 180))
    out = io.BytesIO()
    image.save(out, {"image/png": "PNG", "image/jpeg": "JPEG", "image/webp": "WEBP"}[mime])
    return Artifact(out.getvalue(), mime, f"sample.{mime.rsplit('/', 1)[1]}")


@pytest.mark.parametrize("mime", ["image/png", "image/jpeg", "image/webp"])
def test_image_formats_round_trip_and_recover_wm_code(identity, profile, mime):
    registry = AdapterRegistry()
    result = registry.for_mime(mime).personalize(image_artifact(mime), identity, profile, wm_code=123456)
    assert result.artifact.data != image_artifact(mime).data
    recovered = registry.for_mime(mime).carrier.detect(result.artifact, profile)
    assert recovered.raw["wmCode"] == 123456


def _metadata_stripped_artifact(image: Image.Image, *, mime: str = "image/png") -> Artifact:
    stream = io.BytesIO()
    image.save(stream, {"image/png": "PNG", "image/jpeg": "JPEG"}[mime], quality=60)
    return Artifact(stream.getvalue(), mime, f"stripped.{mime.rsplit('/', 1)[1]}")


def test_image_visual_code_recovers_after_metadata_loss_jpeg_and_resize(identity, profile):
    """The raster channel, not native metadata, carries this recovery."""
    source = Image.new("RGB", (512, 512), (120, 150, 180))
    stream = io.BytesIO(); source.save(stream, "PNG")
    registry = AdapterRegistry()
    personalized = registry.for_mime("image/png").personalize(
        Artifact(stream.getvalue(), "image/png", "large.png"), identity, profile, wm_code=0xAABBCCDD
    )
    raster = Image.open(io.BytesIO(personalized.artifact.data)).convert("RGB")
    jpeg = registry.for_mime("image/jpeg").carrier.detect(_metadata_stripped_artifact(raster, mime="image/jpeg"), profile)
    resized = raster.resize((384, 384), Image.Resampling.LANCZOS)
    resized_evidence = registry.for_mime("image/png").carrier.detect(_metadata_stripped_artifact(resized), profile)
    for evidence in (jpeg, resized_evidence):
        assert evidence.raw["recovery"] == "visual-raster"
        assert evidence.raw["wmCode"] == 0xAABBCCDD
        assert evidence.raw["crcValid"] is True
        assert evidence.raw["syncScore"] >= evidence.raw["acceptanceThresholds"]["minSyncScore"]


def test_image_visual_code_has_explicit_center_crop_threshold_and_random_negative(identity, profile):
    source = Image.new("RGB", (512, 512), (120, 150, 180))
    stream = io.BytesIO(); source.save(stream, "PNG")
    carrier = AdapterRegistry().for_mime("image/png").carrier
    personalized = carrier.embed(Artifact(stream.getvalue(), "image/png", "large.png"), identity, profile, wm_code=0xAABBCCDD)
    raster = Image.open(io.BytesIO(personalized.artifact.data)).convert("RGB")
    recovered = carrier.detect(_metadata_stripped_artifact(raster.crop((128, 128, 384, 384))), profile)
    insufficient = carrier.detect(_metadata_stripped_artifact(raster.crop((166, 166, 345, 345))), profile)
    random = Image.fromarray(np.fromfunction(lambda y, x, c: (x * 17 + y * 31 + c * 71) % 256, (512, 512, 3), dtype=int).astype(np.uint8))
    negative = carrier.detect(_metadata_stripped_artifact(random), profile)
    assert recovered.raw["recovery"] == "visual-raster"
    assert recovered.raw["wmCode"] == 0xAABBCCDD
    assert insufficient.raw["recovery"] == "none"
    assert negative.raw["recovery"] == "none"
    assert "acceptanceThresholds" in negative.raw


def test_pdf_personalization_is_valid_and_text_remains(identity, profile):
    source = fitz.open()
    page = source.new_page()
    page.insert_text((72, 72), "A selectable document")
    original = Artifact(source.tobytes(), "application/pdf", "sample.pdf")
    result = AdapterRegistry().for_mime(original.mime_type).personalize(original, identity, profile)
    checked = fitz.open(stream=result.artifact.data, filetype="pdf")
    assert "selectable document" in checked[0].get_text()
    assert result.artifact.data != original.data


def test_docx_and_pptx_remain_openable(identity, profile):
    doc = Document()
    doc.add_paragraph("Native document")
    stream = io.BytesIO(); doc.save(stream)
    doc_artifact = Artifact(stream.getvalue(), "application/vnd.openxmlformats-officedocument.wordprocessingml.document", "sample.docx")
    doc_result = AdapterRegistry().for_mime(doc_artifact.mime_type).personalize(doc_artifact, identity, profile)
    assert Document(io.BytesIO(doc_result.artifact.data)).paragraphs[0].text == "Native document"

    deck = Presentation()
    deck.slides.add_slide(deck.slide_layouts[6])
    stream = io.BytesIO(); deck.save(stream)
    ppt_artifact = Artifact(stream.getvalue(), "application/vnd.openxmlformats-officedocument.presentationml.presentation", "sample.pptx")
    ppt_result = AdapterRegistry().for_mime(ppt_artifact.mime_type).personalize(ppt_artifact, identity, profile)
    assert len(Presentation(io.BytesIO(ppt_result.artifact.data)).slides) == 1


def test_unsupported_input_has_typed_error():
    with pytest.raises(UnsupportedFormatError) as raised:
        AdapterRegistry().for_mime("text/plain")
    assert raised.value.code == "UNSUPPORTED_FORMAT"
